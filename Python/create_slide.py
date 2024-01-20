# -*- coding: utf-8 -*-
import numpy as np
#define basic chart
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.util import Pt
#data label
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import MSO_AUTO_SIZE

from datetime import datetime
from dateutil.parser import parse

import time
from import_data import input_data

def create_slide(titles):
    slide=prs.slides.add_slide(prs.slide_layouts[5])
    title=slide.shapes.title
    title.text=titles
    return slide

def create_chart(inputdata,slide,chart_type,values,columns,index,data_label,chart_location):
    chart_type_dict={"bar":XL_CHART_TYPE.BAR_CLUSTERED,
                    "bar stacked":XL_CHART_TYPE.BAR_STACKED,
                    "bar stacked 100":XL_CHART_TYPE.BAR_STACKED_100,
                    "column":XL_CHART_TYPE.COLUMN_CLUSTERED,
                    "column stacked":XL_CHART_TYPE.COLUMN_STACKED,
                    "column stacked 100":XL_CHART_TYPE.COLUMN_STACKED_100,
                    "line":XL_CHART_TYPE.LINE_MARKERS}  
    
    location={4.1:[0.2,1],4.2:[4.7,1],4.3:[0.2,4],4.4:[4.7,4],3.1:[2.5,1],3.2:[0.2,4],3.3:[4.7,4]}

    
    base_data=inputdata.pivot_table(values,columns,index,fill_value=0,aggfunc=np.sum)
    base_data.style.format("{:.2%}")
    list_category=list(base_data.index)
    base_data=dict(base_data)
    
    chart_data=ChartData()
    chart_data.categories=list_category
    keys=list(base_data.keys())
    for n in range(len(base_data)):
        chart_data.add_series(keys[n],base_data[keys[n]])
        
    #chart_data.add_series()   
    x,y,cx,cy=Inches(location[chart_location][0]),Inches(location[chart_location][1]),Inches(5.5),Inches(3)
    
    chart=slide.shapes.add_chart(chart_type_dict[chart_type],x,y,cx,cy,chart_data).chart
    
    chart.has_title=False
    chart.series[0].smooth = True
    value_axis = chart.value_axis
    tick_labels = value_axis.tick_labels
    tick_labels.number_format="0,0"
    tick_labels.font.size=Pt(12)
    
    #chart legend            
    chart.has_legend=True
    chart.legend.position=XL_LEGEND_POSITION.RIGHT
    chart.legend.font.size=Pt(14)
    chart.legend.font.name="Arial"
    chart.legend.include_in_layout=False
    
    plot=chart.plots[0]
    if data_label==True:
        plot.has_data_labels=True
        data_labels=plot.data_labels
        
        data_labels.font.size=Pt(12)
        data_labels.font.color.rgb=RGBColor(128,128,128)
        data_labels.number_format="0,50"
        if chart_type=="bar" or chart_type=="column":
            data_labels.position=XL_LABEL_POSITION.OUTSIDE_END

def create_pie(inputdata,slide,values,index,chart_location):

    location={4.1:[0.2,1],4.2:[4.7,1],4.3:[0.2,4],4.4:[4.7,4],3.1:[2.5,1],3.2:[0.2,4],3.3:[4.7,4]}
    
    base_data=inputdata.pivot_table(values=values,index=index,aggfunc=np.sum)
    total_sum=sum(base_data[values])
    base_data["PERCENT"]=base_data[values]/total_sum
             
    chart_data=ChartData()
    chart_data.categories=list(base_data.index.values)
    chart_data.add_series("Percent",list(base_data["PERCENT"]))
    
    x,y,cx,cy=Inches(location[chart_location][0]),Inches(location[chart_location][1]),Inches(5.5),Inches(3)
    chart=slide.shapes.add_chart(XL_CHART_TYPE.PIE, x,y,cx,cy,chart_data).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size=Pt(15)

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

def create_table(inputdata,slide,values,index,columns,table_location):
    base_data=inputdata.pivot_table(values=values,index=index,columns=columns,aggfunc=np.sum)
    base_array=np.array(base_data)
    #slide=prs.slides.add_slide(prs.slide_layouts[5])
    shapes=slide.shapes
    
    location={2.1:[1,1],2.2:[1,4]}
    
    column_sum=[]
    for n in range(len(base_array)):
        column_sum.append(sum(base_array[n]))
    index_sum=[]
    for n in range(len(base_array.T)):
        index_sum.append(sum(base_array.T[n]))
    if sum(column_sum)!=sum(index_sum):
        print("not done")
    
    rows=base_data.shape[0]+2
    cols=base_data.shape[1]+2
    top=Inches(location[table_location][0])
    left=Inches(location[table_location][1])
    width = Inches(8.5)
    height = Inches(1.0)
    table = shapes.add_table(rows, cols, top,left, width, height).table
    
    table_column=list(base_data.columns)
    table_column.append("Total")
    table_column.append("Total")
    table_index=list(base_data.index)
    table_index.append("Total")
    table_index.append("Total")
    for row in range(rows):
        for col in range(cols):
                #table.cell(row,col).text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            if row==0 and col>0:
                table.cell(row,col).text=str(table_column[col-1])                
                #table.cell(row,col).text_frame.fit_text()
            if col==0 and row>0:
                table.cell(row,col).text=str(table_index[row-1])
            if row==0 and col==0:
                table.cell(row,col).text=""
            table.cell(row,col).text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            if row!=0 and col!=0:                
                try: 
                    table.cell(row,col).text=str(base_array[row-1,col-1])
                except IndexError:
                    if row==(rows-1) and col>1:
                        table.cell(row,col-1).text=str(index_sum[col-2])
                    if col==(cols-1) and row>1:
                        table.cell(row-1,col).text=str(column_sum[row-2])
                    else:
                        pass
            if row==(rows-1) and col==(cols-1):
                table.cell(row,col).text=str(sum(column_sum))
            #table.cell(row,col).text_frame.fit_text(max_size=18)
#4 charts: [0.2,1],[4.7,1],[0.2,4],[4.7,4]#
prs=Presentation()
start=time.time()
slide_1=create_slide(titles="combo chart")
inputdata=input_data[(input_data['region-base'].isin(['Roanoke','West']))&(input_data['month'].between(7,9))]

print(time.time()-start)
prs.save("D:/test.pptx")
